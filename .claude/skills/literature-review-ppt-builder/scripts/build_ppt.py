#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Literature Review PPT Builder

å°† ppt_content.md è½¬æ¢ä¸º .pptx æ¼”ç¤ºæ–‡ç¨¿ï¼Œä½¿ç”¨æŒ‡å®šæ¨¡æ¿ã€‚

Usage:
    python build_ppt.py <ppt_content.md> <figures_dir> <template.pptx> <output.pptx>

Example:
    python build_ppt.py output/ppt_content.md output/ Template.pptx output/presentation.pptx
"""

import os
import re
import sys
import copy
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image


# ============================================
# é…ç½®å¸¸é‡
# ============================================

# å¹»ç¯ç‰‡å°ºå¯¸ (16:9 æ ‡å‡†)
SLIDE_WIDTH_INCH = 13.333  # æ ‡å‡† 16:9 å®½åº¦
SLIDE_HEIGHT_INCH = 7.5    # æ ‡å‡† 16:9 é«˜åº¦

# å¸ƒå±€é…ç½®
LAYOUT_CONFIG = {
    # ç¬¬ä¸€é¡µï¼ˆå«æ ‡é¢˜ï¼‰ï¼šä» 3/10 å¼€å§‹
    "first_slide": {
        "content_start_ratio": 0.30,  # 3/10
    },
    # åç»­é¡µï¼šä» 1/5 å¼€å§‹
    "normal_slide": {
        "content_start_ratio": 0.20,  # 1/5
    },
    # å·¦å³åˆ†å‰²ï¼š3:2
    "left_ratio": 0.60,   # å·¦ä¾§æ–‡å­—åŒºåŸŸå  60%
    "right_ratio": 0.40,  # å³ä¾§å›¾ç‰‡åŒºåŸŸå  40%
    # è¾¹è·
    "margin_left": 0.04,   # å·¦è¾¹è· 4%
    "margin_right": 0.04,  # å³è¾¹è· 4%
    "gap": 0.02,           # æ–‡å­—å’Œå›¾ç‰‡ä¹‹é—´é—´éš™ 2%
}

# å­—ä½“é…ç½®
FONT_CONFIG = {
    "family": "Times New Roman",
    "title_size": Pt(16),      # è®ºæ–‡æ ‡é¢˜
    "heading1_size": Pt(14),   # ä¸€çº§å¤§çº²
    "heading2_size": Pt(14),   # äºŒçº§å¤§çº²
    "body_size": Pt(12),       # å…¶ä»–å†…å®¹
    "source_size": Pt(12),     # æ¥æºä¿¡æ¯
}

# é¢œè‰²é…ç½®
COLORS = {
    "title": RGBColor(0x1E, 0x3A, 0x5F),      # æ·±è“
    "heading1": RGBColor(0x1E, 0x3A, 0x5F),   # æ·±è“
    "heading2": RGBColor(0x33, 0x33, 0x33),   # æ·±ç°
    "body": RGBColor(0x44, 0x44, 0x44),       # ç°è‰²
    "source": RGBColor(0x66, 0x66, 0x66),     # æµ…ç°
}


# ============================================
# æ•°æ®ç»“æ„
# ============================================

@dataclass
class BulletItem:
    """å¤§çº²é¡¹"""
    level: int  # 1=ä¸€çº§, 2=äºŒçº§, 3=ä¸‰çº§
    text: str
    is_bold: bool = False


@dataclass
class SlideContent:
    """å¹»ç¯ç‰‡å†…å®¹"""
    bullets: List[BulletItem] = field(default_factory=list)
    figures: List[str] = field(default_factory=list)  # å›¾ç‰‡åç§°åˆ—è¡¨
    figure_modes: List[bool] = field(default_factory=list)  # True=content-only
    notes: str = ""


@dataclass
class PresentationData:
    """æ¼”ç¤ºæ–‡ç¨¿æ•°æ®"""
    title: str = ""
    source: str = ""
    slides: List[SlideContent] = field(default_factory=list)


# ============================================
# è§£æ ppt_content.md
# ============================================

def parse_markdown(content: str) -> PresentationData:
    """è§£æ ppt_content.md æ–‡ä»¶"""
    lines = content.split("\n")
    result = PresentationData()
    
    current_slide: Optional[SlideContent] = None
    in_notes = False
    notes_buffer = []
    in_figures = False
    
    for line in lines:
        trimmed = line.strip()
        
        # è§£ææ ‡é¢˜ï¼ˆç¬¬ä¸€ä¸ª # å¼€å¤´ï¼‰
        if not result.title and trimmed.startswith("# "):
            result.title = trimmed[2:].strip()
            continue
        
        # è§£ææ¥æºï¼ˆ*...* æ ¼å¼ï¼‰
        if not result.source and trimmed.startswith("*") and trimmed.endswith("*"):
            result.source = trimmed[1:-1].strip()
            continue
        
        # æ–°å¹»ç¯ç‰‡å¼€å§‹
        if trimmed.startswith("## Slide"):
            if current_slide:
                if notes_buffer:
                    current_slide.notes = "\n".join(notes_buffer).strip()
                result.slides.append(current_slide)
            current_slide = SlideContent()
            in_notes = False
            in_figures = False
            notes_buffer = []
            continue
        
        if not current_slide:
            continue
        
        # è§£æé…å›¾ï¼ˆå¤šè¡Œæ ¼å¼ï¼‰
        if trimmed.startswith("**é…å›¾**:") or trimmed.startswith("**Figures**:"):
            in_figures = True
            in_notes = False
            # æ£€æŸ¥åŒä¸€è¡Œæ˜¯å¦æœ‰å†…å®¹
            rest = trimmed.split(":", 1)[1].strip() if ":" in trimmed else ""
            if rest and not rest.startswith("-"):
                # æ—§æ ¼å¼ï¼šå•è¡Œé€—å·åˆ†éš”
                parse_figure_line(rest, current_slide)
                in_figures = False
            continue
        
        # è§£æé…å›¾åˆ—è¡¨é¡¹
        if in_figures and trimmed.startswith("- "):
            fig_text = trimmed[2:].strip()
            # ç§»é™¤æ³¨é‡Š
            if "#" in fig_text:
                fig_text = fig_text.split("#")[0].strip()
            parse_figure_line(fig_text, current_slide)
            continue
        
        # é…å›¾åŒºåŸŸç»“æŸ
        if in_figures and trimmed and not trimmed.startswith("-"):
            in_figures = False
        
        # è§£æè®²ç¨¿
        if trimmed.startswith("**è®²ç¨¿**:") or trimmed.startswith("**Notes**:"):
            in_notes = True
            in_figures = False
            rest = trimmed.split(":", 1)[1].strip() if ":" in trimmed else ""
            if rest:
                notes_buffer.append(rest)
            continue
        
        # æ”¶é›†è®²ç¨¿å†…å®¹
        if in_notes and not trimmed.startswith("---"):
            notes_buffer.append(line)
            continue
        
        # åˆ†éš”ç¬¦
        if trimmed == "---":
            in_notes = False
            in_figures = False
            continue
        
        # è§£æä¸€çº§å¤§çº² â–¶
        if trimmed.startswith("â–¶"):
            in_figures = False
            # æå–æ–‡æœ¬ï¼Œç§»é™¤ **æ•°å­—. å’Œç»“å°¾çš„ **
            text = re.sub(r"^â–¶\s*\*\*[\d.]+\s*", "", trimmed)
            text = re.sub(r"\*\*$", "", text).strip()
            current_slide.bullets.append(BulletItem(level=1, text=text, is_bold=True))
            continue
        
        # è§£æäºŒçº§å¤§çº² â–¢
        if trimmed.startswith("â–¢"):
            in_figures = False
            text = re.sub(r"^â–¢\s*[\d.]+\s*", "", trimmed).strip()
            current_slide.bullets.append(BulletItem(level=2, text=text))
            continue
        
        # è§£æä¸‰çº§å†…å®¹ï¼ˆ- å¼€å¤´ï¼Œä½†ä¸åœ¨é…å›¾åŒºåŸŸï¼‰
        if trimmed.startswith("- ") and not in_figures and current_slide.bullets:
            text = trimmed[2:].strip()
            current_slide.bullets.append(BulletItem(level=3, text=text))
            continue
    
    # æ·»åŠ æœ€åä¸€ä¸ªå¹»ç¯ç‰‡
    if current_slide:
        if notes_buffer:
            current_slide.notes = "\n".join(notes_buffer).strip()
        result.slides.append(current_slide)
    
    return result


def parse_figure_line(text: str, slide: SlideContent):
    """è§£æå•ä¸ªå›¾ç‰‡å¼•ç”¨"""
    # å¤„ç†é€—å·åˆ†éš”çš„å¤šä¸ªå›¾ç‰‡
    parts = [p.strip() for p in text.split(",")]
    for part in parts:
        if not part:
            continue
        # æ£€æŸ¥æ˜¯å¦æœ‰ [content-only] æ ‡è®°
        content_only = "[content-only]" in part.lower()
        # æå–å›¾ç‰‡åç§°
        fig_name = re.sub(r"\s*\[content-only\]", "", part, flags=re.IGNORECASE).strip()
        if fig_name:
            slide.figures.append(fig_name)
            slide.figure_modes.append(content_only)


# ============================================
# å›¾ç‰‡å¤„ç†
# ============================================

def find_figure_file(figures_dir: str, figure_name: str) -> Optional[str]:
    """æŸ¥æ‰¾å›¾ç‰‡æ–‡ä»¶
    
    æ”¯æŒå¤šç§å‘½åæ ¼å¼ä»¥å…¼å®¹ extract-pdf-figure çš„è¾“å‡ºï¼š
    - Figure_1.png (ç©ºæ ¼è½¬ä¸‹åˆ’çº¿)
    - <pdf_name>_Figure_1.png (å¸¦ PDF åç§°å‰ç¼€)
    - figure_1.png (å°å†™)
    """
    if not figures_dir or not os.path.isdir(figures_dir):
        return None
    
    # æ ‡å‡†åŒ–åç§°: "Figure 1" -> "Figure_1", "Figure 1(a)" -> "Figure_1(a)"
    normalized = figure_name.replace(" ", "_")
    
    # æ„å»ºæ‰€æœ‰å¯èƒ½çš„æ–‡ä»¶åæ¨¡å¼
    possible_names = [
        # ç›´æ¥åŒ¹é…
        normalized + ".png",
        normalized + ".jpg",
        normalized + ".jpeg",
        # å°å†™
        normalized.lower() + ".png",
        normalized.lower() + ".jpg",
        # æ— æ‹¬å·ç‰ˆæœ¬ Figure_1a
        normalized.replace("(", "").replace(")", "") + ".png",
    ]
    
    # å…ˆå°è¯•ç²¾ç¡®åŒ¹é…
    for name in possible_names:
        filepath = os.path.join(figures_dir, name)
        if os.path.exists(filepath):
            return filepath
    
    # å¦‚æœç²¾ç¡®åŒ¹é…å¤±è´¥ï¼Œæœç´¢ç›®å½•ä¸­åŒ…å«å›¾ç‰‡åç§°çš„æ–‡ä»¶
    # è¿™æ ·å¯ä»¥åŒ¹é… "<pdf_name>_Figure_1.png" æ ¼å¼
    try:
        for filename in os.listdir(figures_dir):
            filename_lower = filename.lower()
            normalized_lower = normalized.lower()
            
            # æ£€æŸ¥æ–‡ä»¶åæ˜¯å¦åŒ…å«å›¾ç‰‡åç§°
            if normalized_lower in filename_lower:
                ext = os.path.splitext(filename)[1].lower()
                if ext in ['.png', '.jpg', '.jpeg']:
                    return os.path.join(figures_dir, filename)
            
            # æ£€æŸ¥æ— ç©ºæ ¼ç‰ˆæœ¬
            normalized_no_space = normalized_lower.replace("_", "")
            filename_no_sep = filename_lower.replace("_", "").replace("-", "")
            if normalized_no_space in filename_no_sep:
                ext = os.path.splitext(filename)[1].lower()
                if ext in ['.png', '.jpg', '.jpeg']:
                    return os.path.join(figures_dir, filename)
    except:
        pass
    
    return None


def get_image_size(image_path: str) -> Tuple[int, int]:
    """è·å–å›¾ç‰‡å°ºå¯¸ (width, height)"""
    with Image.open(image_path) as img:
        return img.size


def calculate_image_layout(
    images: List[Tuple[str, int, int]],  # [(path, width, height), ...]
    area_width: float,   # å¯ç”¨åŒºåŸŸå®½åº¦ (inches)
    area_height: float,  # å¯ç”¨åŒºåŸŸé«˜åº¦ (inches)
) -> List[Tuple[str, float, float, float, float]]:
    """
    è®¡ç®—å¤šå¼ å›¾ç‰‡çš„å¸ƒå±€
    
    Returns:
        [(path, x, y, w, h), ...] ç›¸å¯¹äºåŒºåŸŸå·¦ä¸Šè§’çš„ä½ç½®å’Œå°ºå¯¸ (inches)
    """
    if not images:
        return []
    
    n = len(images)
    gap = 0.1  # å›¾ç‰‡é—´éš™ (inches)
    
    if n == 1:
        # å•å¼ å›¾ç‰‡ï¼šå±…ä¸­æ˜¾ç¤º
        path, img_w, img_h = images[0]
        ratio = img_w / img_h
        
        # è®¡ç®—æœ€ä½³å°ºå¯¸
        if ratio > area_width / area_height:
            # å®½å›¾ï¼šä»¥å®½åº¦ä¸ºå‡†
            w = area_width
            h = w / ratio
        else:
            # é«˜å›¾ï¼šä»¥é«˜åº¦ä¸ºå‡†
            h = area_height
            w = h * ratio
        
        x = (area_width - w) / 2
        y = (area_height - h) / 2
        return [(path, x, y, w, h)]
    
    elif n == 2:
        # ä¸¤å¼ å›¾ç‰‡ï¼šæ ¹æ®å®½é«˜æ¯”å†³å®šæ¨ªæ’è¿˜æ˜¯ç«–æ’
        path1, w1, h1 = images[0]
        path2, w2, h2 = images[1]
        ratio1 = w1 / h1
        ratio2 = w2 / h2
        avg_ratio = (ratio1 + ratio2) / 2
        
        if avg_ratio > 1.2:
            # åå®½çš„å›¾ï¼šç«–æ’ï¼ˆä¸Šä¸‹ï¼‰
            cell_h = (area_height - gap) / 2
            results = []
            for i, (path, img_w, img_h) in enumerate(images):
                ratio = img_w / img_h
                if ratio > area_width / cell_h:
                    w = area_width
                    h = w / ratio
                else:
                    h = cell_h
                    w = h * ratio
                x = (area_width - w) / 2
                y = i * (cell_h + gap) + (cell_h - h) / 2
                results.append((path, x, y, w, h))
            return results
        else:
            # åé«˜çš„å›¾ï¼šæ¨ªæ’ï¼ˆå·¦å³ï¼‰
            cell_w = (area_width - gap) / 2
            results = []
            for i, (path, img_w, img_h) in enumerate(images):
                ratio = img_w / img_h
                if ratio > cell_w / area_height:
                    w = cell_w
                    h = w / ratio
                else:
                    h = area_height
                    w = h * ratio
                x = i * (cell_w + gap) + (cell_w - w) / 2
                y = (area_height - h) / 2
                results.append((path, x, y, w, h))
            return results
    
    else:
        # 3-4å¼ å›¾ç‰‡ï¼š2x2 ç½‘æ ¼
        cols = 2
        rows = (n + 1) // 2
        cell_w = (area_width - gap * (cols - 1)) / cols
        cell_h = (area_height - gap * (rows - 1)) / rows
        
        results = []
        for i, (path, img_w, img_h) in enumerate(images):
            row = i // cols
            col = i % cols
            ratio = img_w / img_h
            
            if ratio > cell_w / cell_h:
                w = cell_w
                h = w / ratio
            else:
                h = cell_h
                w = h * ratio
            
            x = col * (cell_w + gap) + (cell_w - w) / 2
            y = row * (cell_h + gap) + (cell_h - h) / 2
            results.append((path, x, y, w, h))
        
        return results


# ============================================
# Markdown æ–‡æœ¬æ ¼å¼å¤„ç†
# ============================================

def add_formatted_text(paragraph, text: str, base_font_name: str, base_font_size, 
                       base_color, base_bold: bool = False):
    """
    è§£æå¹¶æ·»åŠ å¸¦æ ¼å¼çš„æ–‡æœ¬åˆ°æ®µè½
    æ”¯æŒ **ç²—ä½“** æ ¼å¼
    """
    # åŒ¹é… **text** æ¨¡å¼
    pattern = r'\*\*(.+?)\*\*'
    last_end = 0
    
    for match in re.finditer(pattern, text):
        # æ·»åŠ åŒ¹é…å‰çš„æ™®é€šæ–‡æœ¬
        if match.start() > last_end:
            normal_text = text[last_end:match.start()]
            if normal_text:
                run = paragraph.add_run()
                run.text = normal_text
                run.font.name = base_font_name
                run.font.size = base_font_size
                run.font.bold = base_bold
                run.font.color.rgb = base_color
        
        # æ·»åŠ ç²—ä½“æ–‡æœ¬
        bold_text = match.group(1)
        run = paragraph.add_run()
        run.text = bold_text
        run.font.name = base_font_name
        run.font.size = base_font_size
        run.font.bold = True  # å¼ºåˆ¶ç²—ä½“
        run.font.color.rgb = base_color
        
        last_end = match.end()
    
    # æ·»åŠ å‰©ä½™çš„æ™®é€šæ–‡æœ¬
    if last_end < len(text):
        remaining = text[last_end:]
        if remaining:
            run = paragraph.add_run()
            run.text = remaining
            run.font.name = base_font_name
            run.font.size = base_font_size
            run.font.bold = base_bold
            run.font.color.rgb = base_color


# ============================================
# å¹»ç¯ç‰‡å¤åˆ¶å·¥å…·
# ============================================

def duplicate_slide(prs, template_slide):
    """
    å¤åˆ¶æ¨¡æ¿å¹»ç¯ç‰‡ï¼ˆä¿ç•™æ‰€æœ‰èƒŒæ™¯å’Œå…ƒç´ ï¼‰
    
    Args:
        prs: Presentation å¯¹è±¡
        template_slide: è¦å¤åˆ¶çš„æ¨¡æ¿å¹»ç¯ç‰‡
    
    Returns:
        æ–°å¤åˆ¶çš„å¹»ç¯ç‰‡
    """
    from copy import deepcopy
    from pptx.util import Inches
    import copy
    
    # ä½¿ç”¨æ¨¡æ¿å¹»ç¯ç‰‡çš„å¸ƒå±€åˆ›å»ºæ–°å¹»ç¯ç‰‡
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # å¤åˆ¶æ¨¡æ¿å¹»ç¯ç‰‡çš„æ‰€æœ‰å½¢çŠ¶
    for shape in template_slide.shapes:
        # è·³è¿‡å ä½ç¬¦ï¼ˆå®ƒä»¬å·²ç»é€šè¿‡å¸ƒå±€ç»§æ‰¿äº†ï¼‰
        if shape.is_placeholder:
            continue
            
        # å¤åˆ¶å›¾ç‰‡
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                # è·å–å›¾ç‰‡æ•°æ®
                image_blob = shape.image.blob
                image_content_type = shape.image.content_type
                
                # æ·»åŠ å›¾ç‰‡åˆ°æ–°å¹»ç¯ç‰‡
                from io import BytesIO
                image_stream = BytesIO(image_blob)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            except Exception as e:
                print(f"  è­¦å‘Š: å¤åˆ¶å›¾ç‰‡å¤±è´¥: {e}")
        
        # å¤åˆ¶ç›´çº¿å’Œè¿æ¥å™¨
        elif shape.shape_type in (9, 21):  # LINE=9, STRAIGHT_CONNECTOR=21
            try:
                # è·å–ç›´çº¿çš„èµ·ç‚¹å’Œç»ˆç‚¹åæ ‡
                connector = new_slide.shapes.add_connector(
                    1,  # MSO_CONNECTOR.STRAIGHT
                    shape.begin_x, shape.begin_y,
                    shape.end_x, shape.end_y
                )
                # å¤åˆ¶çº¿æ¡æ ·å¼
                if shape.line.color.type is not None:
                    try:
                        connector.line.color.rgb = shape.line.color.rgb
                    except:
                        pass
                if shape.line.width:
                    connector.line.width = shape.line.width
            except Exception as e:
                # å¦‚æœè¿æ¥å™¨æ–¹å¼å¤±è´¥ï¼Œå°è¯•ç”¨è‡ªé€‰å½¢çŠ¶æ–¹å¼
                try:
                    from pptx.enum.shapes import MSO_SHAPE
                    # è®¡ç®—çº¿æ¡çš„ä½ç½®å’Œå°ºå¯¸
                    left = min(shape.left, shape.left + shape.width)
                    top = min(shape.top, shape.top + shape.height)
                    width = abs(shape.width) if shape.width else shape.left
                    height = abs(shape.height) if shape.height else 1
                    
                    # æ·»åŠ ä¸€ä¸ªçŸ©å½¢ä½œä¸ºæ›¿ä»£çº¿æ¡
                    line_shape = new_slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        shape.left, shape.top,
                        shape.width, shape.height
                    )
                    # è®¾ç½®ä¸ºæ— å¡«å……ã€æœ‰è¾¹æ¡†
                    line_shape.fill.background()
                    if shape.line.width:
                        line_shape.line.width = shape.line.width
                except Exception as e2:
                    print(f"  è­¦å‘Š: å¤åˆ¶ç›´çº¿å¤±è´¥: {e}, {e2}")
        
        # å¤åˆ¶è‡ªé€‰å½¢çŠ¶ï¼ˆçŸ©å½¢ã€åœ†å½¢ç­‰ï¼‰
        elif shape.shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
            try:
                from pptx.enum.shapes import MSO_SHAPE
                new_shape = new_slide.shapes.add_shape(
                    shape.auto_shape_type,
                    shape.left, shape.top,
                    shape.width, shape.height
                )
                # å¤åˆ¶å¡«å……
                if shape.fill.type is not None:
                    try:
                        if shape.fill.type == 1:  # SOLID
                            new_shape.fill.solid()
                            new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                        elif shape.fill.type == 0:  # BACKGROUND
                            new_shape.fill.background()
                    except:
                        pass
                # å¤åˆ¶çº¿æ¡
                try:
                    if shape.line.color.type is not None:
                        new_shape.line.color.rgb = shape.line.color.rgb
                    if shape.line.width:
                        new_shape.line.width = shape.line.width
                except:
                    pass
            except Exception as e:
                print(f"  è­¦å‘Š: å¤åˆ¶å½¢çŠ¶å¤±è´¥: {e}")
        
        # å¤åˆ¶æ–‡æœ¬æ¡†
        elif shape.has_text_frame:
            try:
                new_shape = new_slide.shapes.add_textbox(
                    shape.left, shape.top,
                    shape.width, shape.height
                )
                # å¤åˆ¶æ–‡æœ¬å†…å®¹
                new_tf = new_shape.text_frame
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    if para_idx == 0:
                        new_para = new_tf.paragraphs[0]
                    else:
                        new_para = new_tf.add_paragraph()
                    
                    new_para.alignment = para.alignment
                    new_para.level = para.level
                    
                    for run in para.runs:
                        new_run = new_para.add_run()
                        new_run.text = run.text
                        if run.font.name:
                            new_run.font.name = run.font.name
                        if run.font.size:
                            new_run.font.size = run.font.size
                        if run.font.bold is not None:
                            new_run.font.bold = run.font.bold
                        if run.font.italic is not None:
                            new_run.font.italic = run.font.italic
                        # å®‰å…¨æ£€æŸ¥é¢œè‰²ç±»å‹
                        try:
                            if run.font.color and run.font.color.type is not None:
                                color_rgb = run.font.color.rgb
                                if color_rgb:
                                    new_run.font.color.rgb = color_rgb
                        except:
                            pass  # å¿½ç•¥æ— æ³•è·å–çš„é¢œè‰²
            except Exception as e:
                print(f"  è­¦å‘Š: å¤åˆ¶æ–‡æœ¬æ¡†å¤±è´¥: {e}")
    
    return new_slide


# ============================================
# PPT ç”Ÿæˆ
# ============================================

def create_presentation(
    data: PresentationData,
    figures_dir: str,
    template_path: str,
    output_path: str
):
    """åˆ›å»ºæ¼”ç¤ºæ–‡ç¨¿"""
    
    # åŠ è½½æ¨¡æ¿
    prs = Presentation(template_path)
    
    # è·å–å¹»ç¯ç‰‡å°ºå¯¸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_width_inch = slide_width.inches
    slide_height_inch = slide_height.inches
    
    print(f"  å¹»ç¯ç‰‡å°ºå¯¸: {slide_width_inch:.2f}\" x {slide_height_inch:.2f}\"")
    
    # è·å–æ¨¡æ¿ä¸­çš„ç‰ˆå¼
    # å‡è®¾ï¼šç¬¬ä¸€å¼ å¹»ç¯ç‰‡æ˜¯å°é¢æ¨¡æ¿ï¼Œç¬¬äºŒå¼ æ˜¯å†…å®¹é¡µæ¨¡æ¿
    if len(prs.slides) < 2:
        print("é”™è¯¯: æ¨¡æ¿è‡³å°‘éœ€è¦2å¼ å¹»ç¯ç‰‡")
        sys.exit(1)
    
    # ä¿å­˜æ¨¡æ¿å¹»ç¯ç‰‡å¼•ç”¨ï¼ˆåœ¨åˆ é™¤å‰ï¼‰
    template_first_slide = prs.slides[0]
    template_normal_slide = prs.slides[1]
    first_slide_layout = template_first_slide.slide_layout
    normal_slide_layout = template_normal_slide.slide_layout
    
    # è®¡ç®—å¸ƒå±€å‚æ•°
    margin_left = LAYOUT_CONFIG["margin_left"] * slide_width_inch
    margin_right = LAYOUT_CONFIG["margin_right"] * slide_width_inch
    gap = LAYOUT_CONFIG["gap"] * slide_width_inch
    
    # åˆ†ç•Œçº¿ä½ç½®ï¼ˆ3:2 åˆ†å‰²ï¼‰
    divider_x = slide_width_inch * LAYOUT_CONFIG["left_ratio"]
    
    # å·¦ä¾§æ–‡å­—åŒºåŸŸ
    text_x = margin_left
    text_width = divider_x - margin_left - gap
    
    # å³ä¾§å›¾ç‰‡åŒºåŸŸ
    image_x = divider_x + gap
    image_width = slide_width_inch - divider_x - margin_right - gap
    
    # åˆ›å»ºæ–°çš„æ¼”ç¤ºæ–‡ç¨¿æ¥ä¿å­˜ç»“æœï¼ˆä»æ¨¡æ¿å¤åˆ¶æ•´ä¸ªæ–‡ä»¶ï¼‰
    # ä½†æ˜¯ç›´æ¥ä¿®æ”¹ç°æœ‰çš„å¹»ç¯ç‰‡
    # ç­–ç•¥ï¼šå…ˆå¤åˆ¶æ‰€æœ‰éœ€è¦çš„æ¨¡æ¿å¹»ç¯ç‰‡ï¼Œç„¶ååˆ é™¤åŸæ¨¡æ¿
    
    # å…ˆä¸ºæ¯ä¸ªå†…å®¹å¹»ç¯ç‰‡å¤åˆ¶å¯¹åº”çš„æ¨¡æ¿
    slides_to_create = len(data.slides)
    
    # å¤åˆ¶æ¨¡æ¿å¹»ç¯ç‰‡
    for slide_idx in range(slides_to_create):
        is_first = (slide_idx == 0)
        template_slide = template_first_slide if is_first else template_normal_slide
        # ä½¿ç”¨ duplicate_slide å¤åˆ¶æ¨¡æ¿
        duplicate_slide(prs, template_slide)
    
    # ç°åœ¨åˆ é™¤åŸå§‹çš„ä¸¤ä¸ªæ¨¡æ¿å¹»ç¯ç‰‡
    # åˆ é™¤æ—¶ä»åå¾€å‰åˆ é™¤ï¼Œé¿å…ç´¢å¼•é—®é¢˜
    for _ in range(2):
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # éå†æ¯ä¸ª Slide å¹¶å¡«å……å†…å®¹
    for slide_idx, slide_data in enumerate(data.slides):
        is_first = (slide_idx == 0)
        slide = prs.slides[slide_idx]
        
        # è®¡ç®—å†…å®¹èµ·å§‹ä½ç½®
        if is_first:
            content_start = LAYOUT_CONFIG["first_slide"]["content_start_ratio"] * slide_height_inch
        else:
            content_start = LAYOUT_CONFIG["normal_slide"]["content_start_ratio"] * slide_height_inch
        
        content_height = slide_height_inch - content_start - 0.3  # åº•éƒ¨ç•™ä¸€ç‚¹è¾¹è·
        
        # === ç¬¬ä¸€é¡µï¼šæ›¿æ¢æ ‡é¢˜å’Œæ¥æºæ–‡æœ¬æ¡† ===
        if is_first:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip().upper()
                    if "TITLE" in text:
                        # æ›¿æ¢æ ‡é¢˜
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = data.title
                        run.font.name = FONT_CONFIG["family"]
                        run.font.size = FONT_CONFIG["title_size"]
                        run.font.bold = True
                        run.font.color.rgb = COLORS["title"]
                        p.alignment = PP_ALIGN.CENTER
                    elif "SOURCE" in text:
                        # æ›¿æ¢æ¥æº
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = data.source
                        run.font.name = FONT_CONFIG["family"]
                        run.font.size = FONT_CONFIG["source_size"]
                        run.font.italic = True
                        run.font.color.rgb = COLORS["source"]
                        p.alignment = PP_ALIGN.CENTER
        
        # === æ·»åŠ å¤§çº²å†…å®¹ ===
        if slide_data.bullets:
            # åˆ›å»ºæ–‡æœ¬æ¡†
            txBox = slide.shapes.add_textbox(
                Inches(text_x),
                Inches(content_start),
                Inches(text_width),
                Inches(content_height)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for i, bullet in enumerate(slide_data.bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # è®¾ç½®ç¼©è¿›å’Œé¡¹ç›®ç¬¦å·
                if bullet.level == 1:
                    p.level = 0
                    bullet_char = "â–¶ "
                    font_size = FONT_CONFIG["heading1_size"]
                    font_color = COLORS["heading1"]
                    is_bold = True
                    space_before = Pt(12) if i > 0 else Pt(0)
                elif bullet.level == 2:
                    p.level = 1
                    bullet_char = "â–¢ "
                    font_size = FONT_CONFIG["heading2_size"]
                    font_color = COLORS["heading2"]
                    is_bold = False
                    space_before = Pt(6)
                else:
                    p.level = 2
                    bullet_char = "â€¢ "
                    font_size = FONT_CONFIG["body_size"]
                    font_color = COLORS["body"]
                    is_bold = False
                    space_before = Pt(3)
                
                p.space_before = space_before
                
                # æ·»åŠ å¸¦æ ¼å¼çš„æ–‡æœ¬ï¼ˆå…ˆæ·»åŠ ç¬¦å·ï¼Œå†æ·»åŠ å†…å®¹ï¼‰
                run = p.add_run()
                run.text = bullet_char
                run.font.name = FONT_CONFIG["family"]
                run.font.size = font_size
                run.font.bold = is_bold
                run.font.color.rgb = font_color
                
                # æ·»åŠ å†…å®¹ï¼ˆæ”¯æŒ **ç²—ä½“** æ ¼å¼ï¼‰
                add_formatted_text(p, bullet.text, FONT_CONFIG["family"], 
                                   font_size, font_color, is_bold)
        
        # === æ·»åŠ é…å›¾ ===
        if slide_data.figures:
            # æ”¶é›†å›¾ç‰‡ä¿¡æ¯
            images_info = []
            for fig_name in slide_data.figures:
                fig_path = find_figure_file(figures_dir, fig_name)
                if fig_path:
                    w, h = get_image_size(fig_path)
                    images_info.append((fig_path, w, h))
                else:
                    print(f"  è­¦å‘Š: æ‰¾ä¸åˆ°å›¾ç‰‡ '{fig_name}'")
            
            if images_info:
                # è®¡ç®—å¸ƒå±€
                layout_result = calculate_image_layout(
                    images_info, image_width, content_height
                )
                
                # æ·»åŠ å›¾ç‰‡
                for path, x, y, w, h in layout_result:
                    slide.shapes.add_picture(
                        path,
                        Inches(image_x + x),
                        Inches(content_start + y),
                        Inches(w),
                        Inches(h)
                    )
        
        # === æ·»åŠ  Speaker Notes ===
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.notes
    
    # ä¿å­˜æ–‡ä»¶
    prs.save(output_path)
    print(f"âœ… PPT ç”ŸæˆæˆåŠŸ: {output_path}")
    print(f"   æ€»é¡µæ•°: {len(data.slides)}")


# ============================================
# ä¸»å‡½æ•°
# ============================================

def main():
    if len(sys.argv) < 5:
        print("Usage: python build_ppt.py <ppt_content.md> <figures_dir> <template.pptx> <output.pptx>")
        print("")
        print("Example:")
        print("  python build_ppt.py output/ppt_content.md output/ Template.pptx output/presentation.pptx")
        sys.exit(1)
    
    content_file = sys.argv[1]
    figures_dir = sys.argv[2]
    template_file = sys.argv[3]
    output_file = sys.argv[4]
    
    # æ£€æŸ¥æ–‡ä»¶
    if not os.path.exists(content_file):
        print(f"é”™è¯¯: æ‰¾ä¸åˆ°å†…å®¹æ–‡ä»¶: {content_file}")
        sys.exit(1)
    
    if not os.path.exists(template_file):
        print(f"é”™è¯¯: æ‰¾ä¸åˆ°æ¨¡æ¿æ–‡ä»¶: {template_file}")
        sys.exit(1)
    
    # è¯»å–å†…å®¹
    print(f"ğŸ“– è¯»å–å†…å®¹æ–‡ä»¶: {content_file}")
    with open(content_file, "r", encoding="utf-8") as f:
        content = f.read()
    
    # è§£æ
    data = parse_markdown(content)
    print(f"ğŸ“Š è§£æç»“æœ:")
    print(f"   æ ‡é¢˜: {data.title}")
    print(f"   æ¥æº: {data.source}")
    print(f"   å¹»ç¯ç‰‡æ•°: {len(data.slides)}")
    for i, slide in enumerate(data.slides):
        print(f"   Slide {i+1}: {len(slide.bullets)} ä¸ªè¦ç‚¹, {len(slide.figures)} å¼ é…å›¾")
    
    # ç”Ÿæˆ PPT
    print(f"ğŸ”¨ ç”Ÿæˆ PPT...")
    create_presentation(data, figures_dir, template_file, output_file)


if __name__ == "__main__":
    main()
